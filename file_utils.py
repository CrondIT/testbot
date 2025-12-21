"""Utility functions for file processing and text extraction."""

import os
import tempfile
import PyPDF2
from docx import Document
import pandas as pd
import xlrd  # noqa # Required for legacy XLS files
import pytesseract
from PIL import Image as PILImage
import fitz  # PyMuPDF

# Additional imports for new file formats
from pptx import Presentation
import olefile
from odf.opendocument import load
from odf.text import P
from odf import table, draw


SUPPORTED_EXTENSIONS = [
    ".pdf",
    ".docx",
    ".txt",
    ".xlsx",
    ".xls",
    ".jpg",
    ".jpeg",
    ".png",
    ".bmp",
    ".tiff",
    ".webp",
    ".doc",  # Old Microsoft Word format
    ".odf",  # OpenDocument Text
    ".ods",  # OpenDocument Spreadsheet
    ".odp",  # OpenDocument Presentation
    ".ppt",  # PowerPoint slides
    ".pptx",  # PowerPoint slides (newer format)
]


def get_file_extension(filename: str) -> str:
    """Get file extension from filename"""
    return os.path.splitext(filename)[1]


async def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    try:
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    except Exception as e:
        raise Exception(f"Error processing PDF: {str(e)}")


async def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    try:
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        raise Exception(f"Error processing DOCX: {str(e)}")


async def extract_text_from_txt(file_path: str) -> str:
    """Extract text from TXT file"""
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()
    except UnicodeDecodeError:
        # Try with different encoding
        with open(file_path, "r", encoding="latin-1") as file:
            return file.read()
    except Exception as e:
        raise Exception(f"Error processing TXT: {str(e)}")


async def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from XLSX file"""
    try:
        df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
        text = ""
        for sheet_name, sheet_df in df.items():
            text += f"Sheet: {sheet_name}\n"
            text += sheet_df.to_string()
            text += "\n\n"
        return text
    except Exception as e:
        raise Exception(f"Error processing XLSX: {str(e)}")


async def extract_text_from_xls(file_path: str) -> str:
    """Extract text from XLS file"""
    try:
        df = pd.read_excel(
            file_path, sheet_name=None, engine="xlrd"
        )  # Read all sheets
        text = ""
        for sheet_name, sheet_df in df.items():
            text += f"Sheet: {sheet_name}\n"
            text += sheet_df.to_string()
            text += "\n\n"
        return text
    except Exception as e:
        raise Exception(f"Error processing XLS: {str(e)}")


async def extract_text_from_image(file_path: str) -> str:
    """Extract text from image file using OCR"""
    try:
        # Open the image file
        image = PILImage.open(file_path)

        # Check if Tesseract is available by running a simple test
        import subprocess
        try:
            subprocess.run(
                ["tesseract", "--version"],
                capture_output=True,
                check=True, timeout=5
                        )
        except (
                subprocess.TimeoutExpired,
                subprocess.CalledProcessError,
                FileNotFoundError
                ):
            return "OCR unavailable. Tesseract not installed or not in PATH."

        # Use pytesseract to extract text from the image
        text = pytesseract.image_to_string(
            image, lang="eng+rus"
        )  # Support English and Russian
        return text
    except Exception as e:
        # Handle case when Tesseract is not installed or not working
        if "tesseract" in str(e).lower() or "not found" in str(e).lower():
            return "OCR unavailable. Tesseract not installed or not in PATH."
        else:
            raise Exception(f"Error performing OCR on image: {str(e)}")


async def extract_text_from_pdf_with_ocr(file_path: str) -> str:
    """Extract text from PDF file with OCR fallback for scanned PDFs"""
    try:
        # First try to extract text directly from PDF
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            has_text = False

            for page in pdf_reader.pages:
                page_text = page.extract_text()
                text += page_text + "\n"

                # Check if the page contains substantial text content
                if page_text.strip():
                    has_text = True

        # If the PDF doesn't have much text (likely scanned), use OCR
        if (
            not has_text or len(text.strip()) < 100
        ):  # Threshold to determine if OCR is needed
            # Use PyMuPDF with OCR
            doc = fitz.open(file_path)
            text = ""

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)

                # Try to extract text first
                page_text = page.get_text()

                # If page text is minimal, try OCR
                if len(page_text.strip()) < 50:  # If less than 50 characters
                    # Convert page to image and apply OCR
                    mat = fitz.Matrix(2.0, 2.0)  # Scale for better OCR quality
                    pix = page.get_pixmap(matrix=mat)

                    # Convert to PIL Image
                    img_data = pix.tobytes("png")
                    with tempfile.NamedTemporaryFile(
                        suffix=".png", delete=False
                    ) as temp_img:
                        temp_img.write(img_data)
                        temp_img_path = temp_img.name

                    try:
                        # Check if Tesseract is available 
                        # by running a simple test
                        import subprocess
                        try:
                            subprocess.run(
                                ["tesseract", "--version"],
                                capture_output=True,
                                check=True,
                                timeout=5
                                )
                            # Use pytesseract to extract text from the image
                            page_text = pytesseract.image_to_string(
                                PILImage.open(temp_img_path), lang="eng+rus"
                            )
                        except (
                                subprocess.TimeoutExpired,
                                subprocess.CalledProcessError,
                                FileNotFoundError
                                ):
                            page_text = (
                                f"Page {page_num + 1}: OCR unavailable. "
                                f"Tesseract not installed or not in PATH."
                                )
                    except OSError:
                        # Handle case when Tesseract is not installed
                        page_text = (
                            f"Page {page_num + 1}: OCR unavailable. "
                            f"Tesseract not installed or not in PATH."
                            )
                    finally:
                        # Clean up temporary image file
                        if os.path.exists(temp_img_path):
                            os.remove(temp_img_path)

                text += page_text + "\n"

            doc.close()

        return text
    except Exception as e:
        # Fallback to just text extraction without OCR
        try:
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                if text.strip():
                    return text
                else:
                    raise Exception(f"Error processing PDF with OCR: {str(e)}")
        except Exception:
            raise Exception(f"Error processing PDF with OCR: {str(e)}")


async def extract_text_from_doc(file_path: str) -> str:
    """Extract text from old Microsoft Word .doc file"""
    try:
        # For .doc files, we use olefile to extract text
        ole = olefile.OleFileIO(file_path)
        # Get all stream names
        streams = ole.listdir()

        text = ""
        for stream in streams:
            stream_name = "/".join(stream)
            if ('worddocument' in stream_name.lower() or
                    '1table' in stream_name.lower()):
                try:
                    stream_content = ole.openstream(stream)
                    data = stream_content.read()
                    # Try to decode the binary content to text
                    try:
                        decoded_text = data.decode('utf-8', errors='ignore')
                        text += decoded_text + "\n"
                    except (UnicodeDecodeError, AttributeError):
                        # If direct decoding doesn't work,
                        # try other encodings or skip
                        pass
                except (OSError, IOError):
                    continue  # Skip if stream can't be opened

        # If olefile approach doesn't work well,
        # we may need antiword or similar
        # For now, if text is too short, we'll return an appropriate message
        if len(text.strip()) < 10:
            # Use a temporary approach or return message
            # that this requires external tool
            text = (
                "Текст из .doc файла извлечен частично. "
                "Для полного извлечения требуется дополнительная обработка."
                )

        return text
    except ImportError:
        raise Exception(
            "olefile library not available. "
            "Please install it using: pip install olefile"
            )
    except Exception as e:
        raise Exception(f"Error processing .doc file: {str(e)}")


async def extract_text_from_odf(file_path: str) -> str:
    """Extract text from OpenDocument Text (.odf, .odt) file"""
    try:
        # Open the ODF document
        doc = load(file_path)
        text = ""

        # Extract text from all paragraphs
        for paragraph in doc.getElementsByType(P):
            # Extract text content from paragraph
            text += paragraph.getTextContent() + "\n"

        return text
    except Exception as e:
        raise Exception(f"Error processing ODF file: {str(e)}")


async def extract_text_from_ods(file_path: str) -> str:
    """Extract text from OpenDocument Spreadsheet (.ods) file"""
    try:
        # Open the ODS document
        doc = load(file_path)
        text = ""

        # Process all spreadsheets in the document
        for sheet in doc.getElementsByType(table.Table):
            text += f"Sheet: {sheet.getAttribute('name')}\n"
            # Process all rows in the sheet
            for row in sheet.getElementsByType(table.TableRow):
                for cell in row.getElementsByType(table.TableCell):
                    cell_text = cell.getTextContent()
                    if cell_text:
                        text += cell_text + "\t"
                text += "\n"
            text += "\n"

        return text
    except Exception as e:
        raise Exception(f"Error processing ODS file: {str(e)}")


async def extract_text_from_odp(file_path: str) -> str:
    """Extract text from OpenDocument Presentation (.odp) file"""
    try:
        # Open the ODP document
        doc = load(file_path)
        text = ""

        # Process all presentation slides
        for slide in doc.getElementsByType(draw.Page):
            text += "Slide:\n"
            # Extract text from slide elements
            for element in slide.getElementsByType(P):
                text += element.getTextContent() + "\n"
            text += "\n"

        return text
    except Exception as e:
        raise Exception(f"Error processing ODP file: {str(e)}")


async def extract_text_from_ppt(file_path: str) -> str:
    """Extract text from old PowerPoint .ppt file"""
    try:
        # For .ppt files we need to use alternative methods
        # olefile can be used to extract some information
        ole = olefile.OleFileIO(file_path)
        text = ""

        # Look for powerpoint related streams
        for stream_name in ole.listdir():
            stream_path = "/".join(stream_name)
            if ('powerpoint' in stream_path.lower() or
                    'ppt' in stream_path.lower()):
                try:
                    stream = ole.openstream(stream_name)
                    content = stream.read()
                    # Try to extract text from the binary content
                    content_text = content.decode('utf-8', errors='ignore')
                    if len(content_text.strip()) > 5:
                        text += content_text + "\n"
                except (OSError, IOError, UnicodeDecodeError):
                    continue

        if len(text.strip()) < 5:
            text = (
                "Текст из .ppt файла извлечен частично. "
                "Для полного извлечения требуется дополнительная обработка."
                    )

        return text
    except ImportError:
        raise Exception(
            "olefile library not available. "
            "Please install it using: pip install olefile"
            )
    except Exception as e:
        raise Exception(f"Error processing PPT file: {str(e)}")


async def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint .pptx file"""
    try:
        # Open the presentation
        presentation = Presentation(file_path)
        text = ""

        # Iterate through each slide
        for i, slide in enumerate(presentation.slides):
            text += f"Slide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"

        return text
    except Exception as e:
        raise Exception(f"Error processing PPTX file: {str(e)}")


async def process_uploaded_file(file_path: str, file_extension: str) -> str:
    """Process uploaded file based on its extension and return extracted text
    """
    if file_extension.lower() == ".pdf":
        return await extract_text_from_pdf_with_ocr(file_path)
    elif file_extension.lower() in [
        ".jpg",
        ".jpeg",
        ".png",
        ".bmp",
        ".tiff",
        ".webp",
    ]:
        return await extract_text_from_image(file_path)
    elif file_extension.lower() == ".docx":
        return await extract_text_from_docx(file_path)
    elif file_extension.lower() == ".txt":
        return await extract_text_from_txt(file_path)
    elif file_extension.lower() == ".xlsx":
        return await extract_text_from_xlsx(file_path)
    elif file_extension.lower() == ".xls":
        return await extract_text_from_xls(file_path)
    elif file_extension.lower() == ".doc":
        return await extract_text_from_doc(file_path)
    elif file_extension.lower() in [".odf", ".odt"]:
        return await extract_text_from_odf(file_path)
    elif file_extension.lower() == ".ods":
        return await extract_text_from_ods(file_path)
    elif file_extension.lower() == ".odp":
        return await extract_text_from_odp(file_path)
    elif file_extension.lower() == ".ppt":
        return await extract_text_from_ppt(file_path)
    elif file_extension.lower() == ".pptx":
        return await extract_text_from_pptx(file_path)
    else:
        raise Exception(f"Unsupported file format: {file_extension}")
